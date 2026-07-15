"""
文档解析服务 - 支持 PDF/PPTX/DOCX/Markdown
"""
import os
from pathlib import Path
from typing import List, Tuple
import config


def parse_document(file_path: str) -> Tuple[str, str]:
    """
    解析文档，返回 (纯文本, 文件类型)

    Args:
        file_path: 文件路径

    Returns:
        (提取的文本内容, 文件类型)
    """
    path = Path(file_path)
    suffix = path.suffix.lower()

    parsers = {
        ".pdf": parse_pdf,
        ".pptx": parse_pptx,
        ".ppt": parse_pptx,
        ".docx": parse_docx,
        ".doc": parse_docx,
        ".md": parse_markdown,
        ".markdown": parse_markdown,
        ".txt": parse_txt,
    }

    parser = parsers.get(suffix)
    if not parser:
        raise ValueError(f"不支持的文件格式: {suffix}")

    text = parser(str(path))
    return text, suffix.lstrip(".")


def parse_pdf(file_path: str) -> str:
    """解析PDF文件"""
    import fitz  # PyMuPDF
    doc = fitz.open(file_path)
    texts = []
    for page in doc:
        text = page.get_text("text")
        if text.strip():
            texts.append(text.strip())
    doc.close()
    return "\n\n".join(texts)


def parse_pptx(file_path: str) -> str:
    """解析PPT/PPTX文件"""
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
            elif shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        slide_texts.append(row_text)
        if slide_texts:
            texts.append(f"[幻灯片 {slide_num}]\n" + "\n".join(slide_texts))
    return "\n\n".join(texts)


def parse_docx(file_path: str) -> str:
    """解析DOCX文件"""
    from docx import Document
    doc = Document(file_path)
    texts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            texts.append(text)
    # 也提取表格内容
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                texts.append(row_text)
    return "\n\n".join(texts)


def parse_markdown(file_path: str) -> str:
    """解析Markdown文件"""
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


def parse_txt(file_path: str) -> str:
    """解析纯文本文件"""
    # 尝试多种编码
    for encoding in ["utf-8", "gbk", "gb2312", "utf-16"]:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    raise ValueError(f"无法解码文件: {file_path}")


def chunk_text(text: str, chunk_size: int = None, overlap: int = None) -> List[str]:
    """
    将长文本分割为固定大小的片段（v2.5：中文句感知 + 表格保护）

    Args:
        text: 原始文本
        chunk_size: 片段大小（字符数），默认 config.CHUNK_SIZE
        overlap: 片段重叠大小，默认 config.CHUNK_OVERLAP

    Returns:
        文本片段列表
    """
    chunk_size = chunk_size or config.CHUNK_SIZE
    overlap = overlap or config.CHUNK_OVERLAP

    if not text.strip():
        return []

    # ── 1. 表格检测：连续 |...| 行整体保护，不拆分 ──
    import re
    table_pattern = re.compile(r'^\|(.+)\|.*\|$')
    paragraphs = []
    current_para = []
    in_table = False

    for line in text.split('\n'):
        stripped = line.strip()
        is_table_line = bool(table_pattern.match(stripped))
        if is_table_line:
            if not in_table:
                # 保存上一段非表格内容
                if current_para and not any(table_pattern.match(p) for p in current_para):
                    paragraphs.append('\n'.join(current_para))
                    current_para = []
                in_table = True
            current_para.append(stripped)
        else:
            if in_table and stripped == '':
                # 表格结束
                paragraphs.append('\n'.join(current_para))
                current_para = []
                in_table = False
            elif in_table:
                # 表格内的非表行（紧跟表头/分隔线），继续
                current_para.append(stripped)
                if not table_pattern.match(stripped):
                    paragraphs.append('\n'.join(current_para))
                    current_para = []
                    in_table = False
            else:
                if stripped:
                    current_para.append(stripped)
                else:
                    # 空行 = 段落边界
                    if current_para:
                        paragraphs.append('\n'.join(current_para))
                        current_para = []

    if current_para:
        paragraphs.append('\n'.join(current_para))

    # ── 2. 语义切分：段落级 → 句子级递归 ──
    # 分隔符优先级：\n\n → 。！？ → ； → ，
    sentence_seps = ['。', '！', '？', '；']

    def smart_split(para_text):
        """对超长段落做中文句级切分"""
        if len(para_text) <= chunk_size:
            return [para_text]
        # 尝试在句子边界切分
        parts = []
        start = 0
        for sep in sentence_seps:
            if len(para_text) - start <= chunk_size:
                break
            pos = para_text.rfind(sep, start, min(start + chunk_size, len(para_text)))
            if pos > start:
                parts.append(para_text[start:pos + 1].strip())
                start = pos + 1
        # 剩余部分
        if start < len(para_text):
            parts.append(para_text[start:].strip())
        return parts if parts else [para_text]

    # ── 3. 组装 chunks ──
    chunks = []
    current_chunk = ""

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        if len(current_chunk) + len(para) + 2 <= chunk_size:
            current_chunk = current_chunk + "\n\n" + para if current_chunk else para
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())

            if len(para) > chunk_size:
                # 超长段落 → 用中文句级切分替代硬切
                sub_parts = smart_split(para)
                for sub in sub_parts:
                    chunks.append(sub.strip())
                current_chunk = ""
            else:
                current_chunk = para

    if current_chunk.strip():
        chunks.append(current_chunk.strip())

    # ── 4. 尾段补上 overlap 上下文 ──
    if overlap > 0 and len(chunks) > 1:
        enhanced = []
        for i, c in enumerate(chunks):
            if i > 0 and chunks[i - 1] and len(c) < chunk_size:
                tail = chunks[i - 1][-(overlap // 2):]
                if tail:
                    c = tail + '\n' + c
            enhanced.append(c)
        chunks = enhanced

    return chunks
